import pretty_errors # TODO REMOVE

import collections
import collections.abc
import json
import copy
from io import BytesIO
import pandas as pd
import openpyxl
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData

# IMPORTANT: MUST NAME SLIDE SHAPES https://www.youtube.com/watch?v=IhES3of_9Nw

EXCEL_FILEPATH = "//ACM4/atlanticfiles/Projects/2023 files/EdgeConneX Colo Mystery Shopping - 23024/Cal_verify_quotes_sterilized.xlsx"
EXCEL_SHEET_NAME = "Raw Data"
HEADER_ROW = 3
TARGET_COLUMN = "Quote Name"
SEARCH_TERM = "Cyrus"

with open("columns.json", "r") as columns_file:
    RELEVANT_COLUMNS = json.load(columns_file)
with open("text.json", "r") as text_file:
    ELEMENT_TO_FSTRING = json.load(text_file)
with open("charts.json", "r") as charts_file:
    CHARTS = json.load(charts_file)

def search_excel_sheet(filepath : str, sheet : str, header_row : int, target_column : str, search_term : str) -> pd.DataFrame:
    # Read the Excel file into a DataFrame
    df = pd.read_excel(filepath, sheet_name=sheet, header=header_row)
    # Get all rows that correspond to the search term
    returned_rows = df[df[target_column].str.contains(search_term)]
    # Keep only the specified columns
    relevant_data = returned_rows[list(RELEVANT_COLUMNS.keys())]
    # Rename the columns as specified
    named_data = relevant_data.rename(columns=RELEVANT_COLUMNS)
    # Convert DataFrame to dictionary
    dict_data = named_data.to_dict(orient='list')
    # Check for columns with a single unique value
    collapse_same_vals = False
    if collapse_same_vals:
        for column_name, values in dict_data.items():
            unique_values = set(values)
            if len(unique_values) == 1:
                dict_data[column_name] = unique_values.pop()

    return dict_data

def find_shape_in_group(group, shape_name):
    for shape in group.shapes:
        if shape.shape_type == 6:  # 6 represents a grouped shape
            found_shape = find_shape_in_group(shape, shape_name)
            if found_shape:
                return found_shape
        elif shape.name == shape_name:
            return shape

def get_shape_by_name(slide, shape_name):
    # TODO elegantly raise exception when shape is missing? Can't just put in recursive method cuz it will never find
    # first check placholders     
    for shape in slide.placeholders:
        if shape.name == shape_name:
            return shape
    # If not found, now check shapes (recursively to check groups)
    return find_shape_in_group(slide, shape_name)

def update_charts(slide, provider_data):
    for chart_name, data_cols in CHARTS.items():
        # Retrieve data for chart as per charts.json
        categories = data_cols
        values = [provider_data.get(data_col) for data_col in data_cols]
        print(categories)
        print(values)
        # Modify chart excel workbook
        chart_obj = get_shape_by_name(slide, chart_name).chart
        before = copy.deepcopy(chart_obj)

        workbook_blob = chart_obj.part.chart_workbook.xlsx_part.blob
        workbook = openpyxl.load_workbook(BytesIO(workbook_blob))
        sheet_name = "Sheet1"
        sheet = workbook[sheet_name]
        # Find the table
        table = None
        for tbl in sheet.tables.values():
            table = tbl
            break
        # Throw error if table missing
        if not table:
            raise Exception("Failed to find table in the worksheet")
        # Resize the table
        table.name = 'ModifiedTable' # TODO delete this maybe?
        table.range = f"A1:{get_column_letter(len(values[0]))}{len(values)}"
        # Update the values
        for i, row in enumerate(sheet[table.range]):
            for j, cell in enumerate(row):
                cell.value = values[i][j]
                print(cell.value)
        print("updated table range")
        # Save workbook back to blob
        stream = BytesIO()
        workbook.save(stream)
        stream.seek(0)
        chart_obj.part.chart_workbook.xlsx_part.blob = stream.getvalue()
        print(before.part.chart_workbook.xlsx_part.blob)
        print()
        print(chart_obj.part.chart_workbook.xlsx_part.blob)
        print("Saved workbook")

        

# TODO Add slide duplication when they resolve this git issue https://github.com/scanny/python-pptx/issues/132
# Until then, workaround is to manually copy/paste the template slide n times
def update_text(slide):
    for element_name, fstring in ELEMENT_TO_FSTRING.items():
        element = get_shape_by_name(slide, element_name)
        if not element.has_text_frame:
            raise Exception(f"Tried to set text for element '{element}' which does not have a text field")
        text = fstring.format(**provider_data)
        #print(text)
        for paragraph in element.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = text

def generate_slide(index, provider):
    print(f"Generating slide for provider '{provider}'")
    provider_data = search_excel_sheet(EXCEL_FILEPATH, EXCEL_SHEET_NAME, HEADER_ROW, TARGET_COLUMN, provider)
    #print(provider_data)

    # Load the template PowerPoint presentation
    pres = Presentation("template_slide.pptx")

    # Get the desired slide from the template
    slide = pres.slides[index]

    # Manipulate slide
    #print("Updating text objects")
    #update_text(slide)

    print("Updating charts")
    update_charts(slide, provider_data)


    pres.save("template_slide.pptx")
    return index


generate_slide(1, "Cyrus")
#update_charts()